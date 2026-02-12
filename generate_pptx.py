"""Generate VIGA Project Summary PowerPoint from docs/ folder.

Layout per document type:
  Run:      Page 1 = overview (2 large images: input+output)
            Page 2+ = round outputs (6 small images per page)
  Analysis: 1+ pages: text page + large image pages (2 per page)
  Summary:  2+ pages: text page(s) + large image pages (2 per page)
"""
import io
import os
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

DOCS = Path("docs")
OUT_FULL = Path("docs/VIGA_Project_Summary_v5.pptx")
OUT_SMALL = Path("docs/VIGA_Project_Summary_v5_small.pptx")

# Max pixels per inch for embedded images (used only in compressed mode)
EMBED_DPI = 150

# Global flag: when True, non-GIF images are downscaled before embedding
COMPRESS = False

# Colors (dark theme)
BG_DARK = RGBColor(0x0D, 0x11, 0x17)
BG_CARD = RGBColor(0x16, 0x1B, 0x22)
TEXT_PRIMARY = RGBColor(0xE6, 0xED, 0xF3)
TEXT_SECONDARY = RGBColor(0x8B, 0x94, 0x9E)
ACCENT_BLUE = RGBColor(0x58, 0xA6, 0xFF)
ACCENT_GREEN = RGBColor(0x3F, 0xB9, 0x50)
ACCENT_PURPLE = RGBColor(0xA3, 0x71, 0xF7)
ACCENT_ORANGE = RGBColor(0xD2, 0x99, 0x22)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ============================================================================
# Helpers
# ============================================================================

def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18,
             color=TEXT_PRIMARY, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Segoe UI"
    p.alignment = align
    return tb


def add_multiline(slide, left, top, width, height, lines, size=12, color=TEXT_PRIMARY):
    """Add multiple lines of text (each line = a paragraph)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(3)
    return tb


def add_rect(slide, left, top, width, height, fill_color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    s.line.fill.background()
    return s


def _shrink_png(full_path, display_w_inches, display_h_inches):
    """Downscale a non-GIF image to fit display size at EMBED_DPI. Returns BytesIO PNG."""
    max_px_w = int(display_w_inches * EMBED_DPI) if display_w_inches else 2000
    max_px_h = int(display_h_inches * EMBED_DPI) if display_h_inches else 2000
    with Image.open(str(full_path)) as img:
        img = img.convert("RGBA") if img.mode == "RGBA" else img.convert("RGB")
        native_w, native_h = img.size
        if native_w > max_px_w or native_h > max_px_h:
            img.thumbnail((max_px_w, max_px_h), Image.LANCZOS)
        buf = io.BytesIO()
        img.save(buf, format="PNG", optimize=True)
        buf.seek(0)
    return buf, native_w, native_h


def add_img(slide, img_path, left, top, width=None, height=None):
    """Add image preserving aspect ratio, fitted within width x height box.

    When COMPRESS is True, non-GIF images are downscaled to EMBED_DPI.
    GIFs are always embedded as-is (animated) regardless of COMPRESS setting.
    """
    full = DOCS / img_path
    if not full.exists():
        return False
    try:
        is_gif = full.suffix.lower() == ".gif"

        if COMPRESS and not is_gif:
            disp_w_in = width / 914400 if width else None
            disp_h_in = height / 914400 if height else None
            source, native_w, native_h = _shrink_png(full, disp_w_in, disp_h_in)
        else:
            # Use original file — read native dimensions for aspect ratio
            with Image.open(str(full)) as img:
                native_w, native_h = img.size
            source = str(full)

        if width and height:
            aspect = native_w / native_h
            box_aspect = width / height
            if aspect >= box_aspect:
                final_w = width
                final_h = int(width / aspect)
            else:
                final_h = height
                final_w = int(height * aspect)
            x_off = (width - final_w) // 2
            y_off = (height - final_h) // 2
            slide.shapes.add_picture(source, left + x_off, top + y_off,
                                     final_w, final_h)
        elif width:
            slide.shapes.add_picture(source, left, top, width=width)
        elif height:
            slide.shapes.add_picture(source, left, top, height=height)
        else:
            slide.shapes.add_picture(source, left, top)
        return True
    except Exception:
        return False


# ============================================================================
# Page builders
# ============================================================================

def make_title_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG_DARK)
    add_text(s, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
             "VIGA", 60, ACCENT_BLUE, True, PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(3.0), Inches(11), Inches(1),
             "Vision-as-Inverse-Graphics Agent", 28, TEXT_PRIMARY, False, PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(4.0), Inches(11), Inches(0.8),
             "Iterative Generate / Render / Verify Loop using VLMs + Blender",
             18, TEXT_SECONDARY, False, PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
             "Project Summary  |  January - February 2026",
             16, TEXT_SECONDARY, False, PP_ALIGN.CENTER)


def make_flow_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG_DARK)
    add_text(s, Inches(0.5), Inches(0.2), Inches(12), Inches(0.6),
             "Pipeline Architecture", 36, ACCENT_BLUE, True)

    # ---- Left Column: Asset Sources ----
    add_text(s, Inches(0.4), Inches(0.85), Inches(4.0), Inches(0.3),
             "Asset Sources", 16, ACCENT_GREEN, True)

    # Target Image
    add_rect(s, Inches(0.4), Inches(1.15), Inches(4.0), Inches(0.5), BG_CARD)
    add_text(s, Inches(0.4), Inches(1.2), Inches(4.0), Inches(0.4),
             "Target Image  (photograph / artwork)", 12, ACCENT_ORANGE, True,
             PP_ALIGN.CENTER)

    # Arrow
    add_text(s, Inches(0.4), Inches(1.6), Inches(4.0), Inches(0.2),
             "\u25bc", 12, TEXT_SECONDARY, False, PP_ALIGN.CENTER)

    # SAM3D
    add_rect(s, Inches(0.4), Inches(1.8), Inches(4.0), Inches(1.3), BG_CARD)
    add_text(s, Inches(0.6), Inches(1.85), Inches(3.6), Inches(0.3),
             "SAM3D  (image \u2192 3D)", 13, ACCENT_PURPLE, True)
    add_multiline(s, Inches(0.6), Inches(2.15), Inches(3.6), Inches(0.9), [
        "1. SAM ViT-H \u2192 per-object binary masks",
        "2. TRELLIS \u2192 3D mesh per masked region",
        "Out: *.glb (vertex colors) + transform JSON",
        "~2 hours for 8 objects (RTX 5080 GPU)",
    ], 10, TEXT_SECONDARY)

    # Arrow
    add_text(s, Inches(0.4), Inches(3.05), Inches(4.0), Inches(0.2),
             "\u25bc  GLBs + transforms \u2192 Generator", 9, TEXT_SECONDARY,
             False, PP_ALIGN.CENTER)

    # Meshy
    add_rect(s, Inches(0.4), Inches(3.25), Inches(4.0), Inches(1.3), BG_CARD)
    add_text(s, Inches(0.6), Inches(3.3), Inches(3.6), Inches(0.3),
             "Meshy  (text \u2192 3D)", 13, ACCENT_ORANGE, True)
    add_multiline(s, Inches(0.6), Inches(3.6), Inches(3.6), Inches(0.9), [
        "1. Fuzzy-match local cache (instant)",
        "2. API fallback: preview \u2192 refine \u2192 GLB",
        "Out: *.glb (UV textured, high quality)",
        "~0s cached  /  ~7 min per object via API",
    ], 10, TEXT_SECONDARY)

    # ---- Right Column: Core Loop ----
    add_text(s, Inches(4.7), Inches(0.85), Inches(8.2), Inches(0.3),
             "Core Loop  (\u226425 rounds)", 16, ACCENT_GREEN, True)

    # Generator Agent
    add_rect(s, Inches(4.7), Inches(1.15), Inches(8.2), Inches(1.35), BG_CARD)
    add_text(s, Inches(4.9), Inches(1.2), Inches(5.0), Inches(0.3),
             "Generator Agent  (VLM: GPT-5)", 13, ACCENT_GREEN, True)
    add_multiline(s, Inches(4.9), Inches(1.5), Inches(7.8), Inches(0.9), [
        "Writes Blender Python: import GLBs, compose scene, lighting + camera",
        "Selects best GLB per object: SAM3D (shape) vs Meshy (texture)",
        "Tools: initialize_plan | get_better_object | execute_and_evaluate | end",
        "Memory: target image + SAM3D paths + transforms + Verifier feedback",
    ], 10, TEXT_SECONDARY)

    # Arrow
    add_text(s, Inches(4.7), Inches(2.45), Inches(8.2), Inches(0.2),
             "\u25bc  Blender Python script", 9, TEXT_SECONDARY, False,
             PP_ALIGN.CENTER)

    # Executor (left half of right column)
    add_rect(s, Inches(4.7), Inches(2.65), Inches(3.9), Inches(1.35), BG_CARD)
    add_text(s, Inches(4.9), Inches(2.7), Inches(3.5), Inches(0.3),
             "Blender Executor  (EEVEE)", 13, ACCENT_BLUE, True)
    add_multiline(s, Inches(4.9), Inches(3.0), Inches(3.5), Inches(0.9), [
        "Headless: blender --background",
        "Wrapper sets EEVEE Next engine",
        "Runs Generator's Python script",
        "Out: Camera*.png + state.blend",
    ], 10, TEXT_SECONDARY)

    # Arrow between executor and verifier
    add_text(s, Inches(8.55), Inches(3.1), Inches(0.5), Inches(0.4),
             "\u2192", 20, TEXT_SECONDARY, False, PP_ALIGN.CENTER)

    # Verifier (right half of right column)
    add_rect(s, Inches(9.0), Inches(2.65), Inches(3.9), Inches(1.35), BG_CARD)
    add_text(s, Inches(9.2), Inches(2.7), Inches(3.5), Inches(0.3),
             "Verifier Agent  (VLM)", 13, ACCENT_BLUE, True)
    add_multiline(s, Inches(9.2), Inches(3.0), Inches(3.5), Inches(0.9), [
        "Compares render \u2194 target image",
        "Checks: layout, materials, scale",
        "Tools: rotate, set_camera, scene_info",
        "Out: structured text feedback",
    ], 10, TEXT_SECONDARY)

    # Feedback loop text
    add_text(s, Inches(4.7), Inches(3.95), Inches(8.2), Inches(0.3),
             "\u21bb  Feedback \u2192 Generator memory \u2192 next round"
             "  (until approved or max rounds)",
             10, ACCENT_PURPLE, False, PP_ALIGN.CENTER)

    # ---- Bottom: Outputs + Modes ----
    # Output Structure
    add_rect(s, Inches(0.4), Inches(4.55), Inches(6.0), Inches(2.7), BG_CARD)
    add_text(s, Inches(0.6), Inches(4.6), Inches(5.6), Inches(0.3),
             "Output Structure", 13, ACCENT_GREEN, True)
    add_multiline(s, Inches(0.6), Inches(4.9), Inches(5.6), Inches(2.2), [
        "output/{mode}/{timestamp}/{task}/",
        "  scripts/{N}.py           Blender Python per round",
        "  renders/{N}/Camera*.png  rendered frames",
        "  renders/{N}/state.blend  scene snapshot (undo)",
        "  generator_memory.json    full LLM conversation",
        "  verifier_memory.json     all feedback rounds",
        "  blender_file.blend       final Blender scene",
        "",
        "data/{mode}/{task}/assets/   Meshy GLB cache",
    ], 10, TEXT_SECONDARY)

    # Pipeline Modes
    add_rect(s, Inches(6.7), Inches(4.55), Inches(6.2), Inches(2.7), BG_CARD)
    add_text(s, Inches(6.9), Inches(4.6), Inches(5.8), Inches(0.3),
             "Pipeline Modes", 13, ACCENT_GREEN, True)
    add_multiline(s, Inches(6.9), Inches(4.9), Inches(5.8), Inches(2.2), [
        "get_asset",
        "  Meshy only \u2014 text-to-3D, no SAM3D",
        "",
        "get_asset_sam3d",
        "  SAM3D auto-init + Meshy replaces bad objects",
        "",
        "Pre-computed (--sam3d-results)",
        "  Load existing SAM3D GLBs, skip reconstruction",
        "",
        "Static: Camera.png  |  Dynamic: Camera_f{NNNN}.png",
    ], 10, TEXT_SECONDARY)


def make_env_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG_DARK)
    add_text(s, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
             "Environment Setup", 36, ACCENT_BLUE, True)
    add_rect(s, Inches(0.5), Inches(1.2), Inches(6.0), Inches(5.8), BG_CARD)
    add_text(s, Inches(0.7), Inches(1.3), Inches(5.6), Inches(0.5),
             "Windows 11 (Primary)", 20, ACCENT_GREEN, True)
    add_multiline(s, Inches(0.7), Inches(1.9), Inches(5.6), Inches(5.0), [
        "OS: Windows 11 Home 10.0.26200",
        "GPU: NVIDIA RTX 5080 16GB VRAM",
        "CPU: AMD Ryzen 9 9900X | 32GB DDR5-6000",
        "",
        "Conda Env: agent (Python 3.10)",
        "  openai 2.6.1  |  mcp 1.20.0",
        "  pillow 12.0  |  numpy 2.2.6",
        "  python-pptx 1.0.2",
        "",
        "Conda Env: sam (Python 3.10)",
        "  torch 2.10.0+cu128 (CUDA 12.8)",
        "  torchvision 0.25.0+cu128",
        "  segment-anything 1.0 (Meta SAM)",
        "  opencv-python 4.13.0",
        "",
        "Blender 4.5",
        "  EEVEE (BLENDER_EEVEE_NEXT)",
        "  Cycles GPU rendering",
    ], 11, TEXT_PRIMARY)
    add_rect(s, Inches(6.8), Inches(1.2), Inches(6.0), Inches(5.8), BG_CARD)
    add_text(s, Inches(7.0), Inches(1.3), Inches(5.6), Inches(0.5),
             "Linux / WSL2 (Tested)", 20, ACCENT_GREEN, True)
    add_multiline(s, Inches(7.0), Inches(1.9), Inches(5.6), Inches(5.0), [
        "WSL2 Ubuntu (SAM3D testing)",
        "  CUDA passthrough from host GPU",
        "  ~30% faster SAM3D inference",
        "  Higher VRAM efficiency",
        "",
        "Windows Compatibility Fixes:",
        "  Path resolution: .resolve() all Paths",
        "  Pipe deadlock: temp files (not pipes)",
        "  Encoding: always utf-8 (not cp1252)",
        "  Spaces in paths: quote subprocess args",
        "",
        "VLM Model: GPT-5 (OpenAI API)",
        "  Generator: writes Blender Python",
        "  Verifier: evaluates renders vs target",
        "",
        "3D Asset Sources:",
        "  SAM3D: image -> 3D mesh (TRELLIS)",
        "  Meshy: text -> 3D mesh (API + cache)",
    ], 11, TEXT_PRIMARY)


def make_text_page(prs, date_str, author, badge, badge_color, title, summary, points):
    """A page with header, summary text, and bullet points."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG_DARK)
    add_text(s, Inches(0.5), Inches(0.2), Inches(9.5), Inches(0.7),
             f"{date_str}  \u2014  {title}", 28, ACCENT_BLUE, True)
    add_text(s, Inches(10.2), Inches(0.3), Inches(2.8), Inches(0.5),
             f"{author}  |  {badge}", 14, badge_color, False, PP_ALIGN.RIGHT)
    y = Inches(1.0)
    if summary:
        add_text(s, Inches(0.5), y, Inches(12.3), Inches(1.1),
                 summary, 16, TEXT_SECONDARY)
        y = Inches(2.2)
    if points:
        add_multiline(s, Inches(0.7), y, Inches(11.9), Inches(7.5) - y - Inches(0.2),
                      points, 16, TEXT_PRIMARY)
    return s


def make_large_image_pages(prs, date_str, title, images):
    """Pages with 2 large images each, side by side, maximized."""
    if not images:
        return
    for i in range(0, len(images), 2):
        batch = images[i:i + 2]
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(s, BG_DARK)
        pg = f" ({i // 2 + 1}/{-(-len(images) // 2)})" if len(images) > 2 else ""
        add_text(s, Inches(0.2), Inches(0.08), Inches(12.9), Inches(0.4),
                 f"{date_str}  \u2014  {title}{pg}", 20, ACCENT_BLUE, True)
        if len(batch) == 1:
            # Single image — center it large
            label, path = batch[0]
            w = Inches(10.0)
            add_text(s, Inches(1.6), Inches(0.5), w, Inches(0.25),
                     label, 11, TEXT_SECONDARY, True, PP_ALIGN.CENTER)
            add_img(s, path, Inches(1.6), Inches(0.8), width=w, height=Inches(6.4))
        else:
            # Two images side by side — maximized
            w = Inches(6.4)
            positions = [Inches(0.15), Inches(6.75)]
            for j, (label, path) in enumerate(batch):
                x = positions[j]
                add_text(s, x, Inches(0.5), w, Inches(0.25),
                         label, 11, TEXT_SECONDARY, True, PP_ALIGN.CENTER)
                add_img(s, path, x, Inches(0.8), width=w, height=Inches(6.4))


def make_run_overview(prs, date_str, author, title, summary, input_img, output_img):
    """Run page 1: overview with input/output images, maximized."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG_DARK)
    add_text(s, Inches(0.2), Inches(0.1), Inches(9.8), Inches(0.5),
             f"{date_str}  \u2014  {title}", 24, ACCENT_BLUE, True)
    add_text(s, Inches(10.2), Inches(0.15), Inches(2.9), Inches(0.4),
             f"{author}  |  Run", 11, ACCENT_GREEN, False, PP_ALIGN.RIGHT)
    add_text(s, Inches(0.2), Inches(0.65), Inches(12.9), Inches(0.55),
             summary, 11, TEXT_SECONDARY)
    img_w = Inches(6.2)
    img_h = Inches(5.5)
    y_label = Inches(1.25)
    y_img = Inches(1.55)
    add_text(s, Inches(0.3), y_label, img_w, Inches(0.25),
             "Input / Target", 12, ACCENT_ORANGE, True, PP_ALIGN.CENTER)
    if input_img:
        add_img(s, input_img, Inches(0.3), y_img, width=img_w, height=img_h)
    add_text(s, Inches(6.8), y_label, img_w, Inches(0.25),
             "Output / Result", 12, ACCENT_GREEN, True, PP_ALIGN.CENTER)
    if output_img:
        add_img(s, output_img, Inches(6.8), y_img, width=img_w, height=img_h)


def make_run_rounds(prs, date_str, title, rounds):
    """Run page 2+: round outputs, 6 per page in 3 cols x 2 rows, maximized."""
    if not rounds:
        return
    per_page = 6
    for idx in range(0, len(rounds), per_page):
        batch = rounds[idx:idx + per_page]
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(s, BG_DARK)
        pg = f" (page {idx // per_page + 1})" if len(rounds) > per_page else ""
        add_text(s, Inches(0.2), Inches(0.08), Inches(12.9), Inches(0.38),
                 f"{date_str}  \u2014  {title}  \u2014  Rounds{pg}", 18, ACCENT_BLUE, True)
        # 3 cols x 2 rows — maximized
        col_w = Inches(4.2)
        img_h = Inches(3.1)
        x_pos = [Inches(0.2), Inches(4.55), Inches(8.9)]
        y_rows = [Inches(0.5), Inches(3.9)]
        for i, (label, path) in enumerate(batch):
            col = i % 3
            row = i // 3
            x = x_pos[col]
            y = y_rows[row]
            add_text(s, x, y, col_w, Inches(0.2),
                     label, 9, TEXT_SECONDARY, True, PP_ALIGN.CENTER)
            add_img(s, path, x, y + Inches(0.22), width=col_w, height=img_h)


def make_closing_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG_DARK)
    add_text(s, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
             "VIGA Project Summary", 48, ACCENT_BLUE, True, PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(3.5), Inches(11), Inches(1),
             "32 docs  |  15 days  |  SAM3D + Meshy + Blender + GPT-5",
             20, TEXT_SECONDARY, False, PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
             "Yuna (win/claude/opus/clawdbot)  |  Arin (wsl/claude/opus/clawdbot)  |  kingyy (win/vscode/opus/hum)  |  Sohee (win/antigravity/gemini-pro-high/clawdbot)",
             16, ACCENT_GREEN, False, PP_ALIGN.CENTER)


# ============================================================================
# Process entry by type
# ============================================================================

def process_entry(prs, date_str, author, entry):
    t = entry["type"]
    if t == "run":
        make_run_overview(prs, date_str, author, entry["title"], entry["summary"],
                          entry.get("input_img"), entry.get("output_img"))
        if entry.get("rounds"):
            make_run_rounds(prs, date_str, entry["title"], entry["rounds"])
    elif t == "analysis":
        make_text_page(prs, date_str, author, "Analysis", ACCENT_PURPLE,
                       entry["title"], entry["summary"], entry.get("key_points", []))
        if entry.get("images"):
            make_large_image_pages(prs, date_str, entry["title"], entry["images"])
    elif t == "summary":
        kp = entry.get("key_points", [])
        imgs = entry.get("images", [])
        if imgs:
            # Page 1: all text, Page 2+: images
            make_text_page(prs, date_str, author, "Summary", ACCENT_ORANGE,
                           entry["title"], entry["summary"], kp)
            make_large_image_pages(prs, date_str, entry["title"], imgs)
        else:
            # No images: split key points across 2 pages
            mid = max(len(kp) // 2, 1)
            make_text_page(prs, date_str, author, "Summary", ACCENT_ORANGE,
                           entry["title"], entry["summary"], kp[:mid])
            make_text_page(prs, date_str, author, "Summary (cont.)", ACCENT_ORANGE,
                           entry["title"], "", kp[mid:] if len(kp) > mid else ["(continued)"])


# ============================================================================
# DATA — organized by date (newest first), each date has entries
# ============================================================================

DATES = [
    # -------------------------------------------------------------------------
    # 2026-02-12
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-12",
        "author": "kingyy (win/vscode/opus/hum)",
        "entries": [
            {
                "type": "analysis",
                "title": "VIGA Pipeline Revisited",
                "summary": "Detailed pipeline architecture: SAM3D reconstruction, Meshy text-to-3D, "
                           "Generator/Verifier agents, Blender executor. Per-module inputs/outputs, "
                           "pipeline modes, environment setup, and agent tools reference.",
                "key_points": [
                    "Pipeline: Target Image -> SAM3D (image->3D) + Meshy (text->3D) -> Generator (Blender Python) -> Executor (EEVEE render) -> Verifier (feedback) -> loop",
                    "SAM3D: SAM segmentation (ViT-H) -> per-object binary masks -> SAM3D/TRELLIS reconstruction -> GLB with vertex colors + transform JSON",
                    "Meshy: local cache fuzzy match (instant) -> API fallback (text-to-3D preview -> refine -> download GLB, ~7 min/object)",
                    "Generator: VLM (GPT-5) writes complete Blender Python scripts per round, calls tools (execute_and_evaluate, get_better_object, initialize_plan)",
                    "Verifier: VLM compares render to target, provides structured text feedback, can inspect scene from multiple angles",
                    "3 pipeline modes: Meshy-only (get_asset), SAM3D+Meshy (get_asset_sam3d), Pre-computed SAM3D (--sam3d-results)",
                    "SAM3D vs Meshy: SAM3D captures scene layout + good for simple shapes; Meshy gives UV-textured meshes for complex objects",
                    "Environment: agent env (Python 3.10, openai, mcp) + sam3d_viga env (Python 3.11, torch 2.5.1+cu121, kaolin 0.17.0)",
                    "Hardware: RTX 5080 16GB, Blender 4.5.5 LTS, Windows 11",
                    "Output per round: scripts/{N}.py + renders/{N}/Camera.png + renders/{N}/state.blend",
                ],
                "images": [],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-11
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-11",
        "author": "kingyy (win/vscode/opus/hum)",
        "entries": [
            {
                "type": "run",
                "title": "Dynamic Scene: Artist Run 2 (SAM3D Fixed)",
                "summary": "SAM3D injection bug fixed. Pre-computed GLBs loaded via --sam3d-results. "
                           "Generator kept 5/8 SAM3D objects (fruits), replaced jug/plate/pear with Meshy cache. "
                           "50% faster than Run 1 (99 min vs 203 min). 21 rounds.",
                "input_img": "test_results_images/dynamic_artist_run2/target.png",
                "output_img": "test_results_images/dynamic_artist_run2/gifs/round_19_rotation.gif",
                "rounds": [
                    ("R2 First SAM3D Import", "test_results_images/dynamic_artist_run2/gifs/round_2.gif"),
                    ("R4 Objects on Table", "test_results_images/dynamic_artist_run2/gifs/round_4.gif"),
                    ("R5 SAM3D+Meshy Mix", "test_results_images/dynamic_artist_run2/gifs/round_5.gif"),
                    ("R15 Best Static", "test_results_images/dynamic_artist_run2/gifs/round_15.gif"),
                    ("R16 Ball Added", "test_results_images/dynamic_artist_run2/gifs/round_16.gif"),
                    ("R17 Ball Impact (360)", "test_results_images/dynamic_artist_run2/gifs/round_17.gif"),
                    ("R19 Rotation (360)", "test_results_images/dynamic_artist_run2/gifs/round_19_rotation.gif"),
                    ("R19 Animation (180f)", "test_results_images/dynamic_artist_run2/gifs/round_19_animation.gif"),
                ],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-10
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-10",
        "author": "kingyy (win/vscode/opus/hum)",
        "entries": [
            {
                "type": "run",
                "title": "Dynamic Scene: Artist Run 1 (SAM3D + Meshy)",
                "summary": "First SAM3D-initialized dynamic scene. 8 objects reconstructed but injection bug "
                           "(generator.py:81) meant Generator never received paths — replaced ALL with Meshy. "
                           "Ball physics too subtle. 18 rounds, 203 min.",
                "input_img": "test_results_images/dynamic_artist_run1/target.png",
                "output_img": "test_results_images/dynamic_artist_run1/gifs/round_18.gif",
                "rounds": [
                    ("SAM3D: Ceramic Jug", "test_results_images/dynamic_artist_run1/sam3d_renders/ceramic_jug.gif"),
                    ("SAM3D: Green Pears", "test_results_images/dynamic_artist_run1/sam3d_renders/green_pears.gif"),
                    ("SAM3D: Orange Pears", "test_results_images/dynamic_artist_run1/sam3d_renders/orange_pears.gif"),
                    ("SAM3D: Plate+Fruits", "test_results_images/dynamic_artist_run1/sam3d_renders/plate_with_fruits.gif"),
                    ("SAM3D: Orange Pear", "test_results_images/dynamic_artist_run1/sam3d_renders/orange_pear.gif"),
                    ("SAM3D: Green Apple", "test_results_images/dynamic_artist_run1/sam3d_renders/green_apple.gif"),
                    ("SAM3D: Green Pear", "test_results_images/dynamic_artist_run1/sam3d_renders/green_pear.gif"),
                    ("SAM3D: Green Apple 2", "test_results_images/dynamic_artist_run1/sam3d_renders/green_apple_1.gif"),
                    ("R2 First Render", "test_results_images/dynamic_artist_run1/gifs/round_2.gif"),
                    ("R9 Scene Refinement", "test_results_images/dynamic_artist_run1/gifs/round_9.gif"),
                    ("R14 Best Static", "test_results_images/dynamic_artist_run1/gifs/round_14.gif"),
                    ("R16 Ball Added (360)", "test_results_images/dynamic_artist_run1/gifs/round_16.gif"),
                    ("R16 Frame 1", "test_results_images/dynamic_artist_run1/round_16/Camera_f0001.png"),
                    ("R16 Frame 320", "test_results_images/dynamic_artist_run1/round_16/Camera_f0320.png"),
                    ("R18 Final (360)", "test_results_images/dynamic_artist_run1/gifs/round_18.gif"),
                    ("R18 Frame 1", "test_results_images/dynamic_artist_run1/round_18/Camera_f0001.png"),
                    ("R18 Frame 260", "test_results_images/dynamic_artist_run1/round_18/Camera_f0260.png"),
                    ("R18 Frame 520", "test_results_images/dynamic_artist_run1/round_18/Camera_f0520.png"),
                ],
            },
            {
                "type": "run",
                "title": "SAM3D+Meshy Static Scene Run 1",
                "summary": "First combined SAM3D+Meshy static scene. SAM3D reconstructed 6 objects but Generator "
                           "replaced all with Meshy (4 cached, 1 new API). 25 rounds, 2h14m.",
                "input_img": "test_results_images/sam3d_meshy_run1/target.png",
                "output_img": "test_results_images/sam3d_meshy_run1/round_17_Cam.png",
                "rounds": [
                    ("R2", "test_results_images/sam3d_meshy_run1/round_02_Camera.png"),
                    ("R5", "test_results_images/sam3d_meshy_run1/round_05_Camera.png"),
                    ("R6", "test_results_images/sam3d_meshy_run1/round_06_Camera.png"),
                    ("R8", "test_results_images/sam3d_meshy_run1/round_08_Cam.png"),
                    ("R10", "test_results_images/sam3d_meshy_run1/round_10_Cam.png"),
                    ("R11", "test_results_images/sam3d_meshy_run1/round_11_Cam.png"),
                    ("R12", "test_results_images/sam3d_meshy_run1/round_12_Cam.png"),
                    ("R13", "test_results_images/sam3d_meshy_run1/round_13_Cam.png"),
                    ("R14", "test_results_images/sam3d_meshy_run1/round_14_Cam.png"),
                    ("R15", "test_results_images/sam3d_meshy_run1/round_15_Cam.png"),
                    ("R16", "test_results_images/sam3d_meshy_run1/round_16_Cam.png"),
                    ("R17 Final", "test_results_images/sam3d_meshy_run1/round_17_Cam.png"),
                ],
            },
            {
                "type": "analysis",
                "title": "SAM3D+Meshy Combined Pipeline Design",
                "summary": "Technical design for combining SAM3D reconstruction with Meshy text-to-3D. "
                           "SAM3D initializes scene from target image, Meshy replaces poor-quality objects.",
                "key_points": [
                    "Combined flow: SAM3D auto-init -> segment target -> reconstruct GLBs",
                    "Generator receives SAM3D object paths in memory context",
                    "Phase 1: Render SAM3D scene, evaluate each object quality",
                    "Phase 1: Call get_better_object (Meshy) for poor objects only",
                    "Phase 2: Compose scene with best GLBs (SAM3D or Meshy mix)",
                    "",
                    "Code changes to generator.py: capture SAM3D results, inject into memory",
                    "New prompt: get_asset_sam3d (three-phase workflow)",
                    "Bug fixes: kaolin DLL error, pipeline.yaml path resolution",
                    "Conda env mapping: sam env for SAM3D, agent env for Generator",
                ],
                "images": [],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-09
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-09",
        "author": "kingyy (win/vscode/opus/hum)",
        "entries": [
            {
                "type": "run",
                "title": "Static Scene: Green Tea Run 11",
                "summary": "19 rounds with AABB collision avoidance and non-uniform table scaling. "
                           "Bottle Japanese label text visible by round 11. Cached Meshy assets.",
                "input_img": "test_results_images/greentea/target.png",
                "output_img": "test_results_images/static_scene/20260209_143141/rotation_gif/round_19.gif",
                "rounds": [
                    ("R1", "test_results_images/static_scene/20260209_143141/rotation_gif/round_1.gif"),
                    ("R3", "test_results_images/static_scene/20260209_143141/rotation_gif/round_3.gif"),
                    ("R13", "test_results_images/static_scene/20260209_143141/rotation_gif/round_13.gif"),
                    ("R19 Final", "test_results_images/static_scene/20260209_143141/rotation_gif/round_19.gif"),
                ],
            },
            {
                "type": "run",
                "title": "Static Scene: Green Tea Run 10",
                "summary": "19 rounds. Night-desk aesthetic with blue monitor glow. Introduced orient_group_min_z() "
                           "rotation helper for finding flattest object orientation. Cached Meshy assets.",
                "input_img": "test_results_images/greentea/target.png",
                "output_img": None,
                "rounds": [],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-08
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-08",
        "author": "kingyy (win/vscode/opus/hum)",
        "entries": [
            {
                "type": "run",
                "title": "Static Scene: Green Tea Run 9 (Meshy API)",
                "summary": "First run with Meshy API text-to-3D for all 5 assets. Dramatically higher quality than "
                           "SAM3D: recognizable PET bottle, RGB keyboard, headphones. Cinematic hero shot.",
                "input_img": "test_results_images/greentea/target.png",
                "output_img": "test_results_images/static_scene/20260208_224146/rotation_gif/round_8.gif",
                "rounds": [
                    ("R3", "test_results_images/static_scene/20260208_224146/renders/round_3_Camera.png"),
                    ("R3 (360)", "test_results_images/static_scene/20260208_224146/rotation_gif/round_3.gif"),
                    ("R4", "test_results_images/static_scene/20260208_224146/renders/round_4_Camera.png"),
                    ("R5", "test_results_images/static_scene/20260208_224146/renders/round_5_Camera.png"),
                    ("R5 (360)", "test_results_images/static_scene/20260208_224146/rotation_gif/round_5.gif"),
                    ("R6", "test_results_images/static_scene/20260208_224146/renders/round_6_Camera.png"),
                    ("R7", "test_results_images/static_scene/20260208_224146/renders/round_7_Camera.png"),
                    ("R8 Final", "test_results_images/static_scene/20260208_224146/renders/round_8_Camera.png"),
                    ("R8 Final (360)", "test_results_images/static_scene/20260208_224146/rotation_gif/round_8.gif"),
                ],
            },
            {
                "type": "run",
                "title": "Static Scene: Green Tea Run 8",
                "summary": "First run with Meshy API key configured (but never called — all 5 assets matched locally). "
                           "Slight regression in R6 (objects lost from camera view).",
                "input_img": "test_results_images/greentea/target.png",
                "output_img": "test_results_images/static_scene/20260208_204903/rotation_gif/round_9.gif",
                "rounds": [
                    ("R2", "test_results_images/static_scene/20260208_204903/renders/round_2_Camera.png"),
                    ("R2 (360)", "test_results_images/static_scene/20260208_204903/rotation_gif/round_2.gif"),
                    ("R4", "test_results_images/static_scene/20260208_204903/renders/round_4_Camera.png"),
                    ("R5", "test_results_images/static_scene/20260208_204903/renders/round_5_Camera.png"),
                    ("R5 (360)", "test_results_images/static_scene/20260208_204903/rotation_gif/round_5.gif"),
                    ("R6 Regression", "test_results_images/static_scene/20260208_204903/renders/round_6_Camera.png"),
                    ("R8", "test_results_images/static_scene/20260208_204903/renders/round_8_Camera.png"),
                    ("R9 Final", "test_results_images/static_scene/20260208_204903/renders/round_9_Camera.png"),
                    ("R9 Final (360)", "test_results_images/static_scene/20260208_204903/rotation_gif/round_9.gif"),
                ],
            },
            {
                "type": "run",
                "title": "Static Scene: Green Tea Run 7 (get_asset_simple)",
                "summary": "New prompt setting forbids procedural geometry — only GLB imports allowed. "
                           "All 5 SAM3D assets loaded and arranged over 9 rounds. Zero procedural violations.",
                "input_img": "test_results_images/greentea/target.png",
                "output_img": "test_results_images/static_scene/20260208_164519/rotation_gif/round_9.gif",
                "rounds": [
                    ("R1", "test_results_images/static_scene/20260208_164519/renders/round_1_Camera.png"),
                    ("R1 (360)", "test_results_images/static_scene/20260208_164519/rotation_gif/round_1.gif"),
                    ("R2", "test_results_images/static_scene/20260208_164519/renders/round_2_Camera.png"),
                    ("R3", "test_results_images/static_scene/20260208_164519/renders/round_3_Camera.png"),
                    ("R3 (360)", "test_results_images/static_scene/20260208_164519/rotation_gif/round_3.gif"),
                    ("R4", "test_results_images/static_scene/20260208_164519/renders/round_4_Camera.png"),
                    ("R5", "test_results_images/static_scene/20260208_164519/renders/round_5_Camera.png"),
                    ("R6", "test_results_images/static_scene/20260208_164519/renders/round_6_Camera.png"),
                    ("R6 (360)", "test_results_images/static_scene/20260208_164519/rotation_gif/round_6.gif"),
                    ("R8", "test_results_images/static_scene/20260208_164519/renders/round_8_Camera.png"),
                    ("R9 Final", "test_results_images/static_scene/20260208_164519/renders/round_9_Camera.png"),
                    ("R9 Final (360)", "test_results_images/static_scene/20260208_164519/rotation_gif/round_9.gif"),
                ],
            },
            {
                "type": "run",
                "title": "Static Scene: Green Tea Run 5 (First GLB Import)",
                "summary": "First run where absolute path fix enabled successful GLB asset imports. "
                           "25 rounds, 14 renders. Progression from overexposed early renders to full scene.",
                "input_img": "test_results_images/greentea/target.png",
                "output_img": "test_results_images/static_scene/20260208_050118/renders/round_20_Camera.png",
                "rounds": [
                    ("R3", "test_results_images/static_scene/20260208_050118/renders/round_3_Camera.png"),
                    ("R5", "test_results_images/static_scene/20260208_050118/renders/round_5_Camera.png"),
                    ("R7", "test_results_images/static_scene/20260208_050118/renders/round_7_Camera.png"),
                    ("R9", "test_results_images/static_scene/20260208_050118/renders/round_9_Camera.png"),
                    ("R11", "test_results_images/static_scene/20260208_050118/renders/round_11_Camera.png"),
                    ("R14", "test_results_images/static_scene/20260208_050118/renders/round_14_Camera.png"),
                    ("R16", "test_results_images/static_scene/20260208_050118/renders/round_16_Camera.png"),
                    ("R17", "test_results_images/static_scene/20260208_050118/renders/round_17_Camera.png"),
                    ("R19", "test_results_images/static_scene/20260208_050118/renders/round_19_Camera.png"),
                    ("R20 Final", "test_results_images/static_scene/20260208_050118/renders/round_20_Camera.png"),
                    ("CYCLES Final", "test_results_images/static_scene/20260208_050118/renders/Camera_cycles.png"),
                ],
            },
            {
                "type": "analysis",
                "title": "Asset Pipeline Root Cause (Run 4)",
                "summary": "Root cause analysis: GLB imports failed silently due to relative path resolution in Blender. "
                           "Fix: meshy_api.py resolves to absolute paths with forward slashes.",
                "key_points": [
                    "Problem: Blender CWD != project root -> relative GLB paths fail silently",
                    "All 5 get_better_object calls returned valid paths but Blender couldn't find files",
                    "Generator fell back to procedural geometry for all objects",
                    "Fix in meshy_api.py: Path(previous_assets_dir).resolve() with forward slashes",
                    "Fix in prompts: instructions to use EXACT absolute file paths",
                    "Verified: Run 5 successfully imports all GLB assets after fix",
                ],
                "images": [
                    ("R4 (GLB import failed)", "test_results_images/static_scene/20260207_225603/renders/round_4_Camera.png"),
                    ("R8 (procedural fallback)", "test_results_images/static_scene/20260207_225603/renders/round_8_Camera.png"),
                    ("R13", "test_results_images/static_scene/20260207_225603/renders/round_13_Camera.png"),
                    ("R16", "test_results_images/static_scene/20260207_225603/renders/round_16_Camera.png"),
                    ("R18", "test_results_images/static_scene/20260207_225603/renders/round_18_Camera.png"),
                    ("R19 Final", "test_results_images/static_scene/20260207_225603/renders/round_19_Camera.png"),
                ],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-07
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-07",
        "author": "kingyy (win/vscode/opus/hum)",
        "entries": [
            {
                "type": "run",
                "title": "Static Scene: Green Tea First Run (58 rounds)",
                "summary": "First static scene run with --prompt-setting=none. 58 scripts, 36 renders, ~4 hours. "
                           "No GLB assets (MeshyAPI init failed). Severe scene drift after round 25.",
                "input_img": "test_results_images/greentea/target.png",
                "output_img": "test_results_images/static_scene/20260207_030656/renders/round_58_Camera.png",
                "rounds": [
                    ("R7", "test_results_images/static_scene/20260207_030656/renders/round_7_Camera.png"),
                    ("R9", "test_results_images/static_scene/20260207_030656/renders/round_9_Camera.png"),
                    ("R13", "test_results_images/static_scene/20260207_030656/renders/round_13_Camera.png"),
                    ("R16", "test_results_images/static_scene/20260207_030656/renders/round_16_Camera.png"),
                    ("R19", "test_results_images/static_scene/20260207_030656/renders/round_19_Camera.png"),
                    ("R23", "test_results_images/static_scene/20260207_030656/renders/round_23_Camera.png"),
                    ("R25", "test_results_images/static_scene/20260207_030656/renders/round_25_Camera.png"),
                    ("R28", "test_results_images/static_scene/20260207_030656/renders/round_28_Camera.png"),
                    ("R30", "test_results_images/static_scene/20260207_030656/renders/round_30_Camera.png"),
                    ("R36", "test_results_images/static_scene/20260207_030656/renders/round_36_Camera.png"),
                    ("R42", "test_results_images/static_scene/20260207_030656/renders/round_42_Camera.png"),
                    ("R44", "test_results_images/static_scene/20260207_030656/renders/round_44_Camera.png"),
                    ("R47", "test_results_images/static_scene/20260207_030656/renders/round_47_Camera.png"),
                    ("R50", "test_results_images/static_scene/20260207_030656/renders/round_50_Camera.png"),
                    ("R52", "test_results_images/static_scene/20260207_030656/renders/round_52_Camera.png"),
                    ("R55", "test_results_images/static_scene/20260207_030656/renders/round_55_Camera.png"),
                    ("R57", "test_results_images/static_scene/20260207_030656/renders/round_57_Camera.png"),
                    ("R58 Final", "test_results_images/static_scene/20260207_030656/renders/round_58_Camera.png"),
                    ("CYCLES Final", "test_results_images/static_scene/20260207_030656/renders/CYCLES_final_Camera.png"),
                ],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-06
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-06",
        "author": "kingyy (win/vscode/opus/hum)",
        "entries": [
            {
                "type": "run",
                "title": "Dynamic Scene: Artist Baseline",
                "summary": "First dynamic scene on Windows with GPT-5. Cezanne still life. Procedural geometry only "
                           "(no SAM3D/Meshy). 12 rounds with physics-based ball-throwing animation.",
                "input_img": "test_results_images/dynamic_artist_baseline/target.png",
                "output_img": "test_results_images/dynamic_artist_baseline/gifs/round_12.gif",
                "rounds": [
                    ("R2", "test_results_images/dynamic_artist_baseline/gifs/round_2.gif"),
                    ("R5", "test_results_images/dynamic_artist_baseline/gifs/round_5.gif"),
                    ("R7", "test_results_images/dynamic_artist_baseline/gifs/round_7.gif"),
                    ("R9 (360)", "test_results_images/dynamic_artist_baseline/gifs/round_9.gif"),
                    ("R9 Frame 1", "test_results_images/dynamic_artist_baseline/keyframes/round_9/Camera_f0001.png"),
                    ("R9 Frame 90", "test_results_images/dynamic_artist_baseline/keyframes/round_9/Camera_f0090.png"),
                    ("R9 Frame 180", "test_results_images/dynamic_artist_baseline/keyframes/round_9/Camera_f0180.png"),
                    ("R11", "test_results_images/dynamic_artist_baseline/gifs/round_11.gif"),
                    ("R12 Final (360)", "test_results_images/dynamic_artist_baseline/gifs/round_12.gif"),
                    ("R12 Frame 1", "test_results_images/dynamic_artist_baseline/keyframes/round_12/Camera_f0001.png"),
                    ("R12 Frame 100", "test_results_images/dynamic_artist_baseline/keyframes/round_12/Camera_f0100.png"),
                    ("R12 Frame 200", "test_results_images/dynamic_artist_baseline/keyframes/round_12/Camera_f0200.png"),
                ],
            },
            {
                "type": "analysis",
                "title": "Windows Compatibility Fixes",
                "summary": "Four major Windows issues identified and fixed when running VIGA dynamic scene pipeline.",
                "key_points": [
                    "Issue 1: Blender path resolution",
                    "  Blender resolves relative paths to its own CWD, not project root",
                    "  Fix: .resolve() all Path objects before passing to Blender",
                    "",
                    "Issue 2: Windows path spaces",
                    "  Blender.exe path has spaces: 'C:\\Program Files\\Blender Foundation\\...'",
                    "  Fix: quote subprocess args containing spaces",
                    "",
                    "Issue 3: Subprocess pipe deadlock",
                    "  capture_output=True deadlocks on Windows with large outputs",
                    "  Fix: use temp files for stdout/stderr instead of pipes",
                    "",
                    "Issue 4: cp1252 encoding",
                    "  Windows defaults to cp1252, but scripts contain UTF-8 characters",
                    "  Fix: always specify encoding='utf-8' for file operations",
                    "",
                    "Files fixed: exec.py, investigator_core.py, generator.py, verifier.py, common.py",
                ],
                "images": [],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-05
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-05",
        "author": "Yuna (win/claude/opus/clawdbot)",
        "entries": [
            {
                "type": "summary",
                "title": "SAM3D Pipeline Complete (5/5 Objects)",
                "summary": "End-to-end SAM + SAM3D pipeline completed. 100% success (5/5 objects) on green tea scene. "
                           "SAM segmentation -> mask extraction -> SAM3D reconstruction -> GLB export. ~2 hrs on RTX 5080.",
                "key_points": [
                    "Full pipeline: target.png -> SAM masks -> SAM3D 3D meshes -> GLB export -> Blender render -> GIF",
                    "5 objects: green tea bottle, Ito En bottle, keyboard, headphones, envelope",
                    "All objects successfully reconstructed with vertex colors",
                    "Both X-axis and Y-axis rotation GIFs generated per object",
                    "WebGL viewer integration for interactive 3D preview",
                    "Total processing: ~2 hours on NVIDIA RTX 5080 (16GB VRAM)",
                ],
                "images": [
                    ("Input Scene", "assets/01_greentea_input.jpg"),
                    ("Green Tea Bottle", "assets/green_tea_bottle_y_rotation.gif"),
                    ("Ito En Bottle", "assets/ito_en_bottle_y_rotation.gif"),
                    ("Alienware Keyboard", "assets/alienware_keyboard_y_rotation.gif"),
                    ("Headphones", "assets/headphones_y_rotation.gif"),
                    ("Envelope", "assets/envelope_y_rotation.gif"),
                ],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-04
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-04",
        "author": "Yuna (win/claude/opus/clawdbot)",
        "entries": [
            {
                "type": "summary",
                "title": "SAM Segmentation Results",
                "summary": "SAM segmented green tea desktop scene: 145 raw masks filtered to 8, identifying 5 objects.",
                "key_points": [
                    "Input: green tea desktop photograph",
                    "SAM produced 145 raw masks, filtered to 8 high-confidence regions",
                    "5 objects identified: green tea bottle, Ito En bottle, Alienware keyboard, headphones, envelope",
                    "Mask quality: clean edges, proper object boundaries",
                    "Ready for SAM3D batch 3D reconstruction",
                ],
                "images": [
                    ("Target Scene", "test_results_images/greentea/target.png"),
                    ("Green Tea Bottle Mask", "test_results_images/sam_init/green_tea_bottle.png"),
                    ("Ito En Bottle Mask", "test_results_images/sam_init/ito_en_bottle.png"),
                    ("Keyboard Mask", "test_results_images/sam_init/alienware_keyboard.png"),
                    ("Headphones Mask", "test_results_images/sam_init/headphones.png"),
                    ("Envelope Mask", "test_results_images/sam_init/envelope.png"),
                ],
            },
            {
                "type": "summary",
                "title": "SAM3D All-Masks Results",
                "summary": "SAM3D all-masks pipeline: 6 objects reconstructed with X and Y rotation GIFs.",
                "key_points": [
                    "6 objects: green tea bottle, Ito En bottle, bottle cap, label wrap, bottle neck, headphones",
                    "Both X-axis and Y-axis rotation GIFs generated for each object",
                    "Unexpected sub-object detection (bottle cap, label wrap, neck)",
                    "GLB generation statistics and mask coverage analysis",
                ],
                "images": [
                    ("Input", "test_results_images/01_greentea_input.jpg"),
                    ("Green Tea Bottle", "test_results_images/all_masks_results/green_tea_bottle_y_rotation.gif"),
                    ("Ito En Bottle", "test_results_images/all_masks_results/ito_en_green_tea_bottle_y_rotation.gif"),
                    ("Bottle Cap", "test_results_images/all_masks_results/bottle_cap_y_rotation.gif"),
                    ("Label Wrap", "test_results_images/all_masks_results/label_wrap_y_rotation.gif"),
                    ("Headphones", "test_results_images/all_masks_results/headphones_y_rotation.gif"),
                ],
            },
            {
                "type": "analysis",
                "title": "SAM3D AllMasks Test Investigation",
                "summary": "4 rounds of mask quality investigation. Direct pipeline test failed (missing open3d). "
                           "Discovered mask format issue and reassessed quality metrics.",
                "key_points": [
                    "Round 1: Direct SAM3D pipeline test failed (missing open3d dependency)",
                    "Round 2: VIGA architecture discovery during debugging",
                    "Round 3: Mask quality analysis with black-and-white visualization",
                    "Round 4: Corrected quality reassessment — initial misjudgment about coverage metrics",
                    "Key finding: mask format must be binary (0/1), not grayscale (0/255)",
                ],
                "images": [
                    ("Mask Comparison", "test_results_images/visualizations/20260204_mask_comparison_blackwhite.png"),
                ],
            },
            {
                "type": "analysis",
                "title": "SAM3D Rotation Optimization",
                "summary": "Quaternion vs Euler rotation for rendering SAM3D GLBs in Blender. "
                           "Quaternion rotation avoids gimbal lock.",
                "key_points": [
                    "Problem: Euler rotation causes gimbal lock at certain angles",
                    "Solution: quaternion rotation for smooth 360-degree renders",
                    "Test object: Ito En green tea bottle",
                    "Generated X-axis and Y-axis rotation GIFs + frame samples",
                    "Quaternion rotation integrated into blender_render_rotation.py",
                ],
                "images": [
                    ("SAM Input", "test_results_images/test_sam/ito_en_green_tea_bottle.png"),
                    ("X-Rotation", "test_results_images/bottle_obj_rotate/ito_en_green_tea_bottle_x_rotation.gif"),
                    ("Y-Rotation", "test_results_images/bottle_obj_rotate/ito_en_green_tea_bottle_y_rotation.gif"),
                    ("X Frame 0", "test_results_images/bottle_obj_rotate/ito_en_green_tea_bottle_x_00.png"),
                    ("Y Frame 0", "test_results_images/bottle_obj_rotate/ito_en_green_tea_bottle_y_00.png"),
                    ("Y Frame 15", "test_results_images/bottle_obj_rotate/ito_en_green_tea_bottle_y_15.png"),
                ],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-03
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-03",
        "author": "Yuna (win/claude/opus/clawdbot)",
        "entries": [
            {
                "type": "run",
                "title": "SAM3D Bottle Test (Windows)",
                "summary": "Successful SAM3D bottle reconstruction. Single green tea bottle -> 379K-vertex 3D mesh "
                           "(15MB GLB) in ~9 minutes on RTX 5080. Proper geometry and vertex colors confirmed.",
                "input_img": "test_results_images/test_sam/ito_en_green_tea_bottle.png",
                "output_img": "test_results_images/ito_en_green_tea_bottle_render.png",
                "rounds": [],
            },
            {
                "type": "run",
                "title": "SAM3D Windows Native Test",
                "summary": "First SAM3D on Windows 11. Single-image to 3D mesh generation. "
                           "Windows outperformed WSL for SAM3D inference speed.",
                "input_img": "test_results_images/01_greentea_input.jpg",
                "output_img": "test_results_images/02_greentea_output.png",
                "rounds": [],
            },
            {
                "type": "analysis",
                "title": "Debug: Flat Mesh Issue",
                "summary": "Root cause: SAM3D produced flat table shape instead of bottle. Wrong input image used "
                           "(inverted mask showing table background with bottle cut out).",
                "key_points": [
                    "Symptom: SAM3D output was a flat table-shaped mesh instead of a bottle",
                    "Investigation: compared correct vs incorrect input images",
                    "Root cause: inverted mask — table background passed instead of bottle foreground",
                    "SAM3D correctly reconstructed what it was given (table shape)",
                    "Fix: use original segmentation mask, not its inverse",
                ],
                "images": [
                    ("Wrong Input (inverted mask)", "test_results_images/test_sam/green_tea_bottle.png"),
                    ("Flat Mesh Result", "test_results_images/green_tea_bottle_viga_render.png"),
                    ("Correct Input", "test_results_images/test_sam/ito_en_green_tea_bottle.png"),
                    ("Correct Result", "test_results_images/ito_en_green_tea_bottle_render.png"),
                ],
            },
            {
                "type": "analysis",
                "title": "VIGA Workflow: SAM3D Call",
                "summary": "Complete SAM3D call from image to 3D mesh. Documents pipeline workflow, output JSON, "
                           "mesh structure (flat 'billboard' mesh problem), and disabled mesh post-processing root cause.",
                "key_points": [
                    "Documented complete SAM3D inference call chain",
                    "Output JSON includes transform data (translation, rotation, scale)",
                    "Mesh analysis revealed flat 'billboard' geometry",
                    "Root cause: mesh post-processing was disabled in pipeline config",
                    "Fix: enable mesh post-processing in pipeline.yaml",
                ],
                "images": [
                    ("Input Scene", "test_results_images/01_greentea_input.jpg"),
                    ("Flat Mesh Render", "test_results_images/green_tea_bottle_viga_render.png"),
                ],
            },
            {
                "type": "analysis",
                "author": "Sohee (win/antigravity/gemini-pro-high/clawdbot)",
                "title": "VIGA Architecture Reference",
                "summary": "Technical reference: VIGA -> SAM3D -> TRELLIS hierarchy, inference pipeline stages, "
                           "dual decoder architecture, VRAM distribution, and optimization configs.",
                "key_points": [
                    "Three-tier hierarchy: VIGA (orchestrator) -> SAM3D (reconstruction) -> TRELLIS (core model)",
                    "VIGA: dual-agent system — Generator writes Blender Python, Verifier evaluates renders",
                    "SAM3D: Meta SAM segmentation -> TRELLIS 3D reconstruction per mask",
                    "TRELLIS: sparse structure generation -> dual decoder (Gaussian Splatting + Mesh)",
                    "VRAM budget: ~14GB peak on RTX 5080 (16GB available)",
                    "Optimization: half precision, gradient checkpointing, batch size tuning",
                ],
                "images": [],
            },
            {
                "type": "analysis",
                "title": "SAM3D Internals Deep Dive",
                "summary": "Technical investigation of SAM3D/TRELLIS pipeline internals. "
                           "4 stages, tunable parameters, and VRAM optimization strategies.",
                "key_points": [
                    "Stage 1 (Preprocess): image resize, normalization, feature extraction",
                    "Stage 2 (Sparse Structure): voxel grid generation, sparse 3D structure",
                    "Stage 3 (Sparse Latent): latent code generation from sparse structure",
                    "Stage 4 (Decode): dual decoder — Gaussian Splatting + Mesh extraction",
                    "",
                    "Tunable params: seed, guidance_strength, sampling_steps",
                    "num_gaussians cannot be changed at runtime (hardcoded in model)",
                    "Decode is the bottleneck (~60% of total inference time)",
                    "VRAM optimization: half precision saves ~3GB, gradient checkpointing saves ~2GB",
                ],
                "images": [],
            },
            {
                "type": "analysis",
                "title": "SAM3D Tools Reference",
                "summary": "CLI tools for the SAM3D workflow: inference, visualization, rendering, and GIF creation.",
                "key_points": [
                    "sam3d_worker.py: main inference script (image + mask -> GLB)",
                    "downsample_pointcloud.py: reduce vertex count for large meshes",
                    "render_pointcloud.py: visualize 3D point clouds in Blender",
                    "blender_render_rotation.py: generate 360-degree rotation renders",
                    "create_gifs.py: assemble rendered frames into animated GIFs",
                    "Complete workflow: segment -> reconstruct -> render -> GIF",
                ],
                "images": [],
            },
            {
                "type": "summary",
                "author": "Arin (wsl/claude/opus/clawdbot)",
                "title": "SAM3D WSL Experience Summary",
                "summary": "Consolidated findings from SAM3D testing on WSL2 Ubuntu.",
                "key_points": [
                    "Small images (<=1024px) work on RTX 5080 16GB VRAM",
                    "Large images (>1024px) cause OOM at Stage 3 decode",
                    "Mesh generation must be skipped (Gaussian-only mode)",
                    "Critical discovery: mask format must be 0/1, NOT 0/255",
                    "WSL ~30% faster than Windows for SAM3D inference",
                    "VRAM efficiency better in WSL (less OS overhead)",
                    "",
                    "Recommendation: use Windows native for most tasks",
                    "WSL useful for batch processing with lower VRAM pressure",
                ],
                "images": [],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-02-01
    # -------------------------------------------------------------------------
    {
        "date": "2026-02-01",
        "author": "Arin (wsl/claude/opus/clawdbot)",
        "entries": [
            {
                "type": "run",
                "title": "SAM3D WSL Testing (3 Attempts)",
                "summary": "SAM3D on WSL2 Ubuntu. Three test attempts (original 4480x6720, stage-1-only, 1080p scaled). "
                           "All failed at Stage 3 decode due to OOM on 16GB VRAM.",
                "input_img": "test_results_images/Image_20260130200053_8_44.png",
                "output_img": "test_results_images/Image_20260130200057_9_44.png",
                "rounds": [],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-01-30
    # -------------------------------------------------------------------------
    {
        "date": "2026-01-30",
        "author": "Yuna (win/claude/opus/clawdbot)",
        "entries": [
            {
                "type": "summary",
                "title": "VIGA Progress: SAM3D Environment Setup",
                "summary": "SAM3D environment setup for RTX 5080 (Blackwell sm_120). Critical dependency discovery.",
                "key_points": [
                    "RTX 5080 (sm_120 Blackwell) requires Kaolin 0.18.0",
                    "PyTorch 2.8.0 minimum for Blackwell GPU support",
                    "Three conda environments tested:",
                    "  sam3d: original (incompatible with Blackwell)",
                    "  sam3d_viga: attempted fix (still incompatible)",
                    "  sam3d_kaolin: working (Kaolin 0.18.0 + PyTorch 2.8.0)",
                    "",
                    "Recommended setup: sam3d_kaolin environment",
                    "Successfully imported all SAM3D dependencies",
                    "Next step: run SAM3D inference on test images",
                ],
                "images": [],
            },
        ],
    },

    # -------------------------------------------------------------------------
    # 2026-01-29
    # -------------------------------------------------------------------------
    {
        "date": "2026-01-29",
        "author": "Yuna (win/claude/opus/clawdbot)",
        "entries": [
            {
                "type": "run",
                "title": "VIGA First Tests (Green Tea + Artist)",
                "summary": "Initial VIGA tests with GPT-4o. Two scenarios: (1) green tea bottle tipping over, "
                           "(2) Cezanne still life destruction. Established dual-agent iterative workflow.",
                "input_img": "test_results_images/01_greentea_input.jpg",
                "output_img": "test_results_images/02_greentea_output.png",
                "rounds": [
                    ("Green Tea Input", "test_results_images/01_greentea_input.jpg"),
                    ("Green Tea Output", "test_results_images/02_greentea_output.png"),
                    ("Artist Input", "test_results_images/03_artist_input.png"),
                    ("Artist Output", "test_results_images/04_artist_output.png"),
                ],
            },
            {
                "type": "analysis",
                "author": "Sohee (win/antigravity/gemini-pro-high/clawdbot)",
                "title": "VIGA Workflow Reference",
                "summary": "Comprehensive workflow reference: 5 pipeline modes, dual-agent architecture, "
                           "supported VLMs, 3D tools, and project structure.",
                "key_points": [
                    "5 pipeline modes: dynamic_scene, static_scene, BlenderGym, BlenderBench, SlideBench",
                    "Dual-agent: Generator (writes Blender Python) + Verifier (evaluates renders)",
                    "MCP tool servers for Blender execution and asset management",
                    "Supported VLMs: GPT-4o, GPT-5, Claude",
                    "3D tools: Meshy (text-to-3D), SAM3D (image-to-3D), Infinigen (procedural)",
                    "Iterative loop: generate -> render -> verify -> feedback -> repeat",
                ],
                "images": [],
            },
        ],
    },
]


# ============================================================================
# Main
# ============================================================================

def _build_pptx():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Section 1: Title + Flow + Environment
    make_title_slide(prs)
    make_flow_slide(prs)
    make_env_slide(prs)

    # Section 2: Per-date reports (newest to oldest), summaries first per date
    type_order = {"summary": 0, "analysis": 1, "run": 2}
    for date_group in DATES:
        date_str = date_group["date"]
        author = date_group["author"]
        entries = sorted(date_group["entries"],
                         key=lambda e: type_order.get(e["type"], 9))
        for entry in entries:
            entry_author = entry.get("author", author)
            process_entry(prs, date_str, entry_author, entry)

    # Closing
    make_closing_slide(prs)
    return prs


def main():
    global COMPRESS

    # Full quality version (original images, animated GIFs)
    COMPRESS = False
    prs = _build_pptx()
    prs.save(str(OUT_FULL))
    print(f"Saved: {OUT_FULL} ({len(prs.slides)} slides)")

    # Compressed version (downscaled PNGs, animated GIFs kept as-is)
    COMPRESS = True
    prs = _build_pptx()
    prs.save(str(OUT_SMALL))
    print(f"Saved: {OUT_SMALL} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
